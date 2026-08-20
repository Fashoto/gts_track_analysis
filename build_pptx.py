import pickle
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

EHA_BLUE = RGBColor(0x00, 0x90, 0xFC)
EHA_GREEN = RGBColor(0xE2, 0xEE, 0x64)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x88, 0x88, 0x88)

with open("/tmp/report_data_v3.pkl", "rb") as f:
    R = pickle.load(f)

TEMPLATE = "/Users/BusayoFashoto/.claude/skills/eha-templates/assets/eHA_General_Slide_Template_A.pptx"
prs = Presentation(TEMPLATE)

xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    prs.part.drop_rel(sldId.get(qn("r:id")))
    xml_slides.remove(sldId)

layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}

def add_slide(layout_name):
    return prs.slides.add_slide(layouts[layout_name])

def set_title(slide, text, size=None):
    ph = slide.placeholders[0]
    ph.text_frame.text = text
    if size:
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)
    return ph

def style_body_para(p, text, size=16, bold=False, color=BLACK, level=0):
    p.text = text
    p.level = level
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

def add_textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Open Sans"
    return box

def add_picture_fit(slide, path, left, top, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    l = left + (max_w - w) // 2
    t = top + (max_h - h) // 2
    slide.shapes.add_picture(path, l, t, width=w, height=h)

def remove_placeholder(slide, idx):
    ph = slide.placeholders[idx]
    ph._element.getparent().remove(ph._element)

def add_table(slide, left, top, width, height, data, col_widths=None, font_size=11):
    rows, cols = len(data), len(data[0])
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = gtable.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Open Sans"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = BLACK
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = EHA_BLUE if r == 0 else (RGBColor(0xF2, 0xF9, 0xFF) if r % 2 == 0 else WHITE)
    return gtable

CHART = "charts"
DATE_STR = f"{R['dates'][0].strftime('%d')}–{R['dates'][-1].strftime('%d %B %Y')}"

# =========================================================== 1. TITLE
s = add_slide("TITLE")
s.placeholders[0].text_frame.text = "GPS Track Analysis"
s.placeholders[1].text_frame.text = f"IBRA II — {len(R['dates'])}-Day Round Review (Refreshed, Day {len(R['dates'])} partial) · {DATE_STR}\nData and GIS Analytics, Data Informatics Department, eHealth Africa"

# =========================================================== 2. SECTION: Overview
s = add_slide("SECTION_HEADER_2")
set_title(s, "Overview")

# =========================================================== 3. At a glance
s = add_slide("TITLE_AND_BODY_3")
set_title(s, "At a Glance")
date_range_short = f"{R['dates'][0].strftime('%d')}–{R['dates'][-1].strftime('%d %b')}"
cards = [
    (f"{R['n_total']:,}", "GPS tracks recorded"),
    (f"{len(R['dates'])} days", f"{date_range_short}, Day {len(R['dates'])} partial"),
    (f"{len(R['state_grp'])}", "campaign states covered"),
    (f"{R['n_team_codes']:,}", "distinct team codes"),
    (f"{R['n_devices_deployed']:,}", "devices deployed (roster)"),
    (f"{R['n_imeis']:,}", f"distinct IMEIs seen ({R['device_gap_vs_deployed']:+,} vs. roster)"),
]
for idx, (big, small) in zip([1, 2, 3], cards[:3]):
    ph = s.placeholders[idx]
    tf = ph.text_frame
    tf.paragraphs[0].text = big
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = EHA_BLUE
    p2 = tf.add_paragraph()
    p2.text = small
    p2.font.size = Pt(14)
add_textbox(s, Inches(0.5), Inches(4.55), Inches(9), Inches(0.8),
    [(f"Plus: {cards[3][0]} {cards[3][1]}   ·   {cards[4][0]} {cards[4][1]}   ·   {cards[5][0]} {cards[5][1]}", 12.5, False, GREY)],
    align=PP_ALIGN.CENTER)

# =========================================================== 4. Shape of the round
s = add_slide("TITLE_AND_BODY")
set_title(s, "Shape of the Round: Daily Track Volume")
ph = s.placeholders[1]
ph.text_frame.paragraphs[0].text = (
    f"Volume peaked on Day 2 ({R['peak_day'].strftime('%d %b')}) at {R['peak_day_count']:,} tracks. "
    f"Day {len(R['dates'])} ({R['dates'][-1].strftime('%d %b')}) is still in progress ({R['daily'].iloc[-1]:,} tracks so far)."
)
ph.text_frame.paragraphs[0].font.size = Pt(13)
add_picture_fit(s, f"{CHART}/daily_trend.png", Inches(0.5), Inches(2.05), Inches(9), Inches(3.2))

# =========================================================== 5. SECTION: Data Quality
s = add_slide("SECTION_HEADER_2")
set_title(s, "Data Quality & Methodology")

# =========================================================== 6. Data quality overview
s = add_slide("TITLE_AND_BODY")
set_title(s, "Six Data-Quality Checks Run Against the Full Dataset")
ph = s.placeholders[1]
ph.text_frame.paragraphs[0].text = (
    f"Rates are of the full {R['n_total']:,}-row refreshed dataset, after CSV repair."
)
ph.text_frame.paragraphs[0].font.size = Pt(13)
add_picture_fit(s, f"{CHART}/qa_overview.png", Inches(0.4), Inches(2.0), Inches(9.2), Inches(3.3))

# =========================================================== 7. QA detail
s = add_slide("TITLE_AND_BODY")
set_title(s, "What Each Check Means")
tf = s.placeholders[1].text_frame
lines = [
    f"Malformed CSV rows ({R['repaired_rows']:,}, {R['repaired_pct']:.2f}%): unescaped commas inside the Team Code field split rows into 16–17 columns instead of 15. Rejoined at source position — no rows dropped.",
    f"Clock mismatch ({R['clock_mismatch_30d']:,}, {R['clock_mismatch_pct']:.2f}%): GPS and phone timestamps disagree by more than 30 days — a device clock fault, not a GPS fault.",
    f"Implausible speed ({R['implausible_speed']:,}, {R['implausible_speed_pct']:.2f}%): fixes above 30 m/s (108 km/h).",
    f"Poor accuracy ({R['poor_accuracy_50m']:,}, {R['poor_accuracy_pct']:.2f}%): reported horizontal accuracy worse than 50m.",
    f"Missing team code ({R['missing_team_code']:,} rows) and unresolved state ({R['state_unmapped_rows']:,} rows, {100-R['state_mapped_pct']:.1f}%): see Team Code Integrity section.",
    f"Zero exact duplicate rows and zero duplicate Track IDs across all {R['n_total']:,} records.",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    style_body_para(p, line, size=13)

# =========================================================== 8. SECTION: Geographic coverage
s = add_slide("SECTION_HEADER_2")
set_title(s, "Geographic Coverage")

# =========================================================== 9. State ranking
s = add_slide("TITLE_AND_BODY")
set_title(s, "Tracks by State")
ph = s.placeholders[1]
ph.text_frame.paragraphs[0].text = (
    f"Kebbi alone accounts for {R['state_grp'].loc['Kebbi','tracks']/R['n_total']*100:.0f}% of all mapped tracks."
)
ph.text_frame.paragraphs[0].font.size = Pt(13)
add_picture_fit(s, f"{CHART}/state_ranking.png", Inches(0.4), Inches(2.0), Inches(9.2), Inches(3.3))

# =========================================================== 10. State coverage inventory
s = add_slide("TITLE_AND_BODY_2")
set_title(s, "State Coverage Inventory")
remove_placeholder(s, 1)
sg = R["state_grp"].reset_index()
data = [["State", "Tracks", "Teams", "Devices", "Tracks/team", "Codes − devices"]]
for _, row in sg.iterrows():
    data.append([row["State"], f"{row['tracks']:,}", f"{row['teams']:,}", f"{row['devices']:,}",
                 f"{row['tracks_per_team']:.0f}", f"{int(row['codes_minus_devices']):+d}"])
add_table(s, Inches(0.3), Inches(1.4), Inches(9.4), Inches(3.9), data, font_size=10.5,
          col_widths=[Inches(1.6), Inches(1.6), Inches(1.3), Inches(1.4), Inches(1.6), Inches(1.9)])

# =========================================================== 10b. Border-area tracks reconciled via team code
s = add_slide("TITLE_AND_BODY_2")
set_title(s, "Border-Area Tracks Reconciled via Team Code")
remove_placeholder(s, 1)
rb = R["oof_recovery_breakdown"]
data = [["Recovered to State", "Tracks"]]
for state, cnt in rb.items():
    data.append([state, f"{cnt:,}"])
add_table(s, Inches(1.6), Inches(1.4), Inches(6.8), Inches(3.4), data, font_size=12.5,
          col_widths=[Inches(4.8), Inches(2.0)])
unresolved_codes = ", ".join(R["oof_unresolved_codes"].index.tolist())
add_textbox(s, Inches(0.6), Inches(4.85), Inches(8.8), Inches(0.7),
    [(f"The polygon join initially placed {R['oof_recovered_by_teamcode']+R['oof_still_unresolved']:,} tracks in a ward outside the 15 campaign states — border-area fieldwork, not GPS error. "
      f"{R['oof_recovered_by_teamcode']:,} were reconciled to their true state via team code; only {R['oof_still_unresolved']:,} remain unresolved (team code \"{unresolved_codes}\" — a test/dummy entry).", 12, False, GREY)])

# =========================================================== 11. Day-one signal — CONFIRMED explanation
s = add_slide("TITLE_AND_BODY")
set_title(s, "Confirmed: Six States Launched on Sunday, Not Saturday")
ph = s.placeholders[1]
states_list = ", ".join(R["sunday_start_states"])
ph.text_frame.paragraphs[0].text = (
    f"{states_list} recorded under 5% of their round's activity on Day 1 (Sat 15 Aug). "
    f"Confirmed operationally: these states began fieldwork on Sunday, 16 Aug — not a data anomaly."
)
ph.text_frame.paragraphs[0].font.size = Pt(12.5)
add_picture_fit(s, f"{CHART}/day1_share.png", Inches(0.4), Inches(2.0), Inches(9.2), Inches(3.3))

# =========================================================== 12. SECTION: Team & Device performance
s = add_slide("SECTION_HEADER_2")
set_title(s, "Team & Device Performance")

# =========================================================== 13. Continuity
s = add_slide("TITLE_AND_BODY")
set_title(s, "Days-Worked Continuity")
cs = R["cont_summary"]
n_complete_days = len(R["dates"]) - 1  # last day is still partial
full_complete = cs.loc[n_complete_days, "track_share_pct"] if n_complete_days in cs.index else 0
ph = s.placeholders[1]
ph.text_frame.paragraphs[0].text = (
    f"{cs.loc[n_complete_days,'count']:,} team codes ({full_complete:.0f}% of all tracks) worked every complete day so far "
    f"(through {R['dates'][n_complete_days-1].strftime('%d %b')}); {cs.loc[1,'count']:,} appear on only 1 day."
)
ph.text_frame.paragraphs[0].font.size = Pt(13)
add_picture_fit(s, f"{CHART}/continuity.png", Inches(0.7), Inches(2.0), Inches(8.6), Inches(3.3))

# =========================================================== 14. Effort distribution
s = add_slide("TITLE_AND_BODY")
set_title(s, "Distribution of Effort Across Teams")
ph = s.placeholders[1]
ph.text_frame.paragraphs[0].text = (
    f"Median team logged {R['team_tracks_median']:.0f} tracks; {R['n_low_vol_teams']:,} codes logged fewer than 50 all round — a supervisor follow-up list."
)
ph.text_frame.paragraphs[0].font.size = Pt(12.5)
add_picture_fit(s, f"{CHART}/effort_hist.png", Inches(0.4), Inches(2.0), Inches(9.2), Inches(3.3))

# =========================================================== 15. Device/code reconciliation (national + by state)
s = add_slide("TITLE_AND_BODY")
set_title(s, "Device Reconciliation vs. Deployment Roster")
tf = s.placeholders[1].text_frame
lines = [
    f"{R['n_devices_deployed']:,} devices were deployed per the roster; {R['n_imeis']:,} distinct IMEIs actually appear in the GPS data — a gap of {R['device_gap_vs_deployed']:+,} ({R['device_gap_vs_deployed']/R['n_devices_deployed']*100:.1f}%).",
    "A positive gap this size points to device swaps mid-round, personal phones used alongside issued devices, or roster records not updated to reflect replacements.",
    f"Separately, {R['imei_multi_code_count']:,} IMEIs ({R['imei_multi_code_pct']:.1f}% of all IMEIs seen) reported under more than one team code across the round — one device reached {R['imei_multi_code_max']} distinct codes.",
    "Recommendation: reconcile the device roster against observed IMEIs before the next round, and confirm whether multi-code devices reflect shared hardware or code re-assignment mid-round.",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    style_body_para(p, line, size=13.5)

# =========================================================== 16. State-level device/code table
s = add_slide("TITLE_AND_BODY_2")
set_title(s, "Device – Team Code Gap by State")
remove_placeholder(s, 1)
sg_sorted = sg.reindex(sg["codes_minus_devices"].abs().sort_values(ascending=False).index)
data = [["State", "Team codes", "Devices", "Codes − devices", "Read"]]
for _, row in sg_sorted.iterrows():
    gap = int(row["codes_minus_devices"])
    if gap > 20:
        read = "Multiple codes reused per device"
    elif gap < -5:
        read = "More devices than codes — check device roster"
    else:
        read = "Near 1:1, healthy"
    data.append([row["State"], f"{row['teams']:,}", f"{row['devices']:,}", f"{gap:+d}", read])
add_table(s, Inches(0.3), Inches(1.4), Inches(9.4), Inches(3.9), data, font_size=11,
          col_widths=[Inches(1.5), Inches(1.4), Inches(1.3), Inches(1.5), Inches(3.7)])

# =========================================================== 17. SECTION: Team Code Integrity
s = add_slide("SECTION_HEADER_2")
set_title(s, "Team Code Integrity")

# =========================================================== 18. Integrity overview
s = add_slide("TITLE_AND_BODY_3")
set_title(s, "Team Code Integrity — at a Glance")
cards2 = [
    (f"{R['state_unmapped_rows']:,}", f"tracks ({100-R['state_mapped_pct']:.1f}%) with an unresolvable state prefix"),
    (f"{R['numeric_only_rows']:,}", "tracks logged under a numeric-only code (e.g. '067')"),
    (f"{R['missing_team_code']:,}", "tracks with a missing/blank team code"),
]
for idx, (big, small) in zip([1, 2, 3], cards2):
    ph = s.placeholders[idx]
    tf = ph.text_frame
    tf.paragraphs[0].text = big
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = EHA_BLUE
    p2 = tf.add_paragraph()
    p2.text = small
    p2.font.size = Pt(12.5)
add_textbox(s, Inches(0.5), Inches(4.55), Inches(9), Inches(0.8),
    [("0 cross-state code collisions found — every code that resolves to a state resolves to exactly one.", 12.5, False, GREY)],
    align=PP_ALIGN.CENTER)

# =========================================================== 19. Highest-volume codes
s = add_slide("TITLE_AND_BODY_2")
set_title(s, "Highest-Volume Team Codes (Round Total)")
remove_placeholder(s, 1)
dup = R["top_duplicated_codes"].reset_index()
dup.columns = ["Team Code", "Tracks"]
data = [["Team Code", "Total tracks"]]
for _, row in dup.iterrows():
    data.append([row["Team Code"], f"{row['Tracks']:,}"])
add_table(s, Inches(1.2), Inches(1.35), Inches(7.6), Inches(4.0), data, font_size=11,
          col_widths=[Inches(5.6), Inches(2.0)])

# =========================================================== 20. Unmapped/malformed detail
s = add_slide("TITLE_AND_BODY")
set_title(s, "What's Behind the Unresolved Codes")
tf = s.placeholders[1].text_frame
lines = [
    f"Numeric-only codes (e.g. 067, 026, 197): serial-style codes with no state/LGA prefix at all — {R['numeric_only_rows']:,} tracks.",
    "Round-name-as-prefix codes (e.g. IBRA/R2/IGB/083): the campaign round name was entered where the state prefix belongs.",
    "Place-name-only codes (e.g. UBANDOMA/IBRA2/FP): a ward/settlement name substituted for the whole code, no state marker.",
    "Ambiguous short prefixes deliberately left unmapped rather than guessed (e.g. KE, BK, KNG) — could plausibly belong to more than one state.",
    f"This {100-R['state_mapped_pct']:.1f}% is a genuine data-collection gap: no dictionary expansion can safely recover a state that was never encoded in the source text.",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    style_body_para(p, line, size=13)

# =========================================================== 21. SECTION: Temporal patterns
s = add_slide("SECTION_HEADER_2")
set_title(s, "Temporal Patterns")

# =========================================================== 22. Teams reporting per day
s = add_slide("TITLE_AND_BODY")
set_title(s, "Teams Reporting Per Day")
ph = s.placeholders[1]
ph.text_frame.paragraphs[0].text = (
    f"Participation rose from {R['teams_per_day'].iloc[0]:,} codes on Day 1 to a peak of {R['teams_per_day'].max():,} on Day 2."
)
ph.text_frame.paragraphs[0].font.size = Pt(13)
add_picture_fit(s, f"{CHART}/teams_per_day.png", Inches(0.6), Inches(2.0), Inches(8.8), Inches(3.3))

# =========================================================== 23. Churn detail
s = add_slide("TITLE_AND_BODY_2")
set_title(s, "Day-to-Day Churn Detail")
remove_placeholder(s, 1)
data = [["Transition", "Retained", "Dropped out", "New codes"]]
for d0, d1, kept, dropped, new in R["churn"]:
    data.append([f"{d0.strftime('%d %b')} → {d1.strftime('%d %b')}", f"{kept:,}", f"{dropped:,}", f"{new:,}"])
add_table(s, Inches(1.0), Inches(1.6), Inches(8.0), Inches(3.2), data, font_size=13,
          col_widths=[Inches(2.6), Inches(1.8), Inches(1.8), Inches(1.8)])
add_textbox(s, Inches(0.6), Inches(4.9), Inches(8.8), Inches(0.5),
    [(f"{R['missed_on_last_complete_day']:,} team codes active before {R['last_complete_day'].strftime('%d %b')} did not report on that day.", 12, False, GREY)])

# =========================================================== 24. Hour-of-day
s = add_slide("TITLE_AND_BODY")
set_title(s, "Hour-of-Day Activity Pattern")
ph = s.placeholders[1]
ph.text_frame.paragraphs[0].text = (
    "Fieldwork is tightly concentrated between 06:00 and 16:00 UTC, consistent with a structured working day."
)
ph.text_frame.paragraphs[0].font.size = Pt(13)
add_picture_fit(s, f"{CHART}/hourly.png", Inches(0.5), Inches(2.0), Inches(9.0), Inches(3.3))

# =========================================================== 25. SECTION: Recommendations
s = add_slide("SECTION_HEADER_2")
set_title(s, "Recommended Actions")

# =========================================================== 26. Recommendations
s = add_slide("TITLE_AND_BODY")
set_title(s, "Recommended Actions")
tf = s.placeholders[1].text_frame
lines = [
    f"Reconcile the device roster ({R['n_devices_deployed']:,}) against the {R['n_imeis']:,} distinct IMEIs observed — a {R['device_gap_vs_deployed']:+,} gap suggests unlogged device swaps or personal-phone use.",
    f"Follow up the {R['n_low_vol_teams']:,} team codes logging under 50 tracks all round, and the {R['missed_on_last_complete_day']:,} codes that stopped reporting on the last complete day.",
    f"Investigate the {R['imei_multi_code_count']:,} devices reporting under more than one team code — confirm shared hardware vs. reassignment.",
    f"Close the {100-R['state_mapped_pct']:.1f}% team-code gap at data-entry time: enforce a state-prefix pick-list on the collection tool rather than free text.",
    f"Fix device clock settings behind the {R['clock_mismatch_30d']:,}-row GPS/phone timestamp mismatch before the next round.",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    style_body_para(p, line, size=13)

# =========================================================== 27. Thank you
s = add_slide("TITLE_1")
s.placeholders[0].text_frame.text = "Thank You"
s.placeholders[1].text_frame.text = "Data and GIS Analytics · Data Informatics Department · eHealth Africa"

OUT = "/Users/BusayoFashoto/Desktop/tracks/export_1200_2026-08-20_12-01-03/Tracks/GPS_Track_Analysis_IBRA_II_Full_Round.pptx"
prs.save(OUT)
print("Saved:", OUT)
